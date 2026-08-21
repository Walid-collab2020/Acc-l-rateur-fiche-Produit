"""Génère le schéma PPTX de présentation client de l'application."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
BLEU_FOND   = RGBColor(0x0D, 0x2B, 0x55)   # fond slide
BLEU_TITRE  = RGBColor(0x1A, 0x3F, 0x7A)   # boites titres
BLEU_CLAIR  = RGBColor(0x22, 0x6E, 0xB0)   # accents
VERT        = RGBColor(0x1E, 0x8C, 0x5A)   # documents sources
ORANGE      = RGBColor(0xE8, 0x76, 0x22)   # IA / traitement
VIOLET      = RGBColor(0x6A, 0x3C, 0x9A)   # fiche générée
ROUGE       = RGBColor(0xC0, 0x39, 0x2B)   # contrôles
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR  = RGBColor(0xEC, 0xF0, 0xF5)
GRIS_TEXTE  = RGBColor(0x2C, 0x3E, 0x50)

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill, line_color=None, line_width=Pt(0), radius=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    fill_elem = shape.fill
    fill_elem.solid()
    fill_elem.fore_color.rgb = fill
    shape.line.fill.background()  # no border by default
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    return shape

def add_rounded_rect(slide, l, t, w, h, fill, line_color=None, line_width=Pt(1.5)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    # Reduce corner radius
    adj = shape.adjustments
    try:
        adj[0] = 0.05
    except Exception:
        pass
    return shape

def set_text(shape, text, font_size=Pt(10), bold=False, color=BLANC, align=PP_ALIGN.CENTER, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color

def add_textbox(slide, text, l, t, w, h, font_size=Pt(9), bold=False, color=BLANC, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_arrow(slide, x1, y1, x2, y2, color=BLEU_CLAIR, width=Pt(2.5)):
    """Flèche droite horizontale ou verticale."""
    from pptx.util import Inches
    connector = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR_TYPE.STRAIGHT if False else 1,  # STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Arrowhead
    from pptx.oxml.ns import qn
    ln = connector.line._ln
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'triangle')
    headEnd.set('w', 'med')
    headEnd.set('len', 'med')


# ── Slide ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank_layout)

# Fond
bg = add_rect(slide, 0, 0, 13.33, 7.5, BLEU_FOND)

# ══════════════════════════════════════════════════════════════════════════════
# TITRE
# ══════════════════════════════════════════════════════════════════════════════
titre_box = add_rect(slide, 0, 0, 13.33, 0.75, BLEU_TITRE)
add_textbox(slide,
    "Application IA de Migration Produit KELIA  —  Vue d'ensemble du fonctionnement",
    0.2, 0.1, 12.9, 0.6,
    font_size=Pt(16), bold=True, color=BLANC, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# COLONNE 1 — DOCUMENTS SOURCES  (x: 0.2 → 3.2)
# ══════════════════════════════════════════════════════════════════════════════
COL1_X = 0.2
COL1_W = 3.0

# En-tête colonne 1
header1 = add_rounded_rect(slide, COL1_X, 0.95, COL1_W, 0.45, VERT)
set_text(header1, "DOCUMENTS SOURCES", Pt(11), bold=True, color=BLANC)

# Sous-label
add_textbox(slide, "Uploadés par produit dans l'application",
    COL1_X, 1.43, COL1_W, 0.3, Pt(8), color=GRIS_CLAIR, align=PP_ALIGN.CENTER)

# Liste des documents
docs = [
    ("📄", "Conditions Générales"),
    ("📄", "Note Technique Actuarielle"),
    ("📄", "Notice d'information"),
    ("📄", "Avenant"),
    ("📊", "Extraction BOSS"),
    ("📝", "CR Atelier / Décision de conception"),
    ("📋", "Fiche Produit existante (si dispo)"),
]

y_doc = 1.78
for icon, label in docs:
    doc_box = add_rounded_rect(slide, COL1_X + 0.08, y_doc, COL1_W - 0.16, 0.42,
                                RGBColor(0x17, 0x5E, 0x3A), line_color=VERT, line_width=Pt(1))
    set_text(doc_box, f"{icon}  {label}", Pt(9), bold=False, color=BLANC, align=PP_ALIGN.LEFT)
    doc_box.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    y_doc += 0.47

# ══════════════════════════════════════════════════════════════════════════════
# COLONNE 2 — TRAITEMENT IA  (x: 3.55 → 7.0)
# ══════════════════════════════════════════════════════════════════════════════
COL2_X = 3.55
COL2_W = 3.4

# Flèche 1 → 2
add_arrow(slide, COL1_X + COL1_W + 0.01, 3.6, COL2_X - 0.01, 3.6, color=ORANGE, width=Pt(3))

# En-tête colonne 2
header2 = add_rounded_rect(slide, COL2_X, 0.95, COL2_W, 0.45, ORANGE)
set_text(header2, "TRAITEMENT PAR L'IA", Pt(11), bold=True, color=BLANC)

add_textbox(slide, "Modèles Claude (Anthropic)",
    COL2_X, 1.43, COL2_W, 0.3, Pt(8), color=GRIS_CLAIR, align=PP_ALIGN.CENTER)

# Étapes de traitement
etapes = [
    ("1", "Classification automatique", "L'IA identifie et classe chaque\ndocument uploadé (type, résumé,\nscore de confiance)"),
    ("2", "Extraction des données produit", "L'IA lit les CG, Notes Techniques,\nAvenants… et en extrait\nles paramètres clés du produit"),
    ("3", "Génération de la Fiche Produit", "L'IA remplit champ par champ\nle template KELIA (.xlsx)\nà partir des données extraites"),
]

y_etape = 1.80
for num, titre_etape, desc in etapes:
    # Boite numéro
    num_box = add_rounded_rect(slide, COL2_X + 0.05, y_etape, 0.38, 0.38,
                                ORANGE, line_color=BLANC, line_width=Pt(1.5))
    set_text(num_box, num, Pt(13), bold=True, color=BLANC)

    # Boite étape
    etape_box = add_rounded_rect(slide, COL2_X + 0.5, y_etape, COL2_W - 0.6, 1.0,
                                  RGBColor(0x6B, 0x3D, 0x00), line_color=ORANGE, line_width=Pt(1))
    tf = etape_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = titre_etape
    r0.font.size = Pt(9)
    r0.font.bold = True
    r0.font.color.rgb = BLANC

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = desc
    r1.font.size = Pt(8)
    r1.font.bold = False
    r1.font.color.rgb = GRIS_CLAIR

    y_etape += 1.12

# ══════════════════════════════════════════════════════════════════════════════
# COLONNE 3 — FICHE GÉNÉRÉE  (x: 7.2 → 9.6)
# ══════════════════════════════════════════════════════════════════════════════
COL3_X = 7.2
COL3_W = 2.4

# Flèche 2 → 3
add_arrow(slide, COL2_X + COL2_W + 0.01, 3.1, COL3_X - 0.01, 3.1, color=VIOLET, width=Pt(3))

# En-tête colonne 3
header3 = add_rounded_rect(slide, COL3_X, 0.95, COL3_W, 0.45, VIOLET)
set_text(header3, "LIVRABLE GÉNÉRÉ", Pt(11), bold=True, color=BLANC)

add_textbox(slide, "Sortie principale",
    COL3_X, 1.43, COL3_W, 0.3, Pt(8), color=GRIS_CLAIR, align=PP_ALIGN.CENTER)

# Fiche KELIA
fiche_box = add_rounded_rect(slide, COL3_X + 0.1, 1.78, COL3_W - 0.2, 1.5,
                               RGBColor(0x3D, 0x1E, 0x6E), line_color=VIOLET, line_width=Pt(1.5))
tf_fiche = fiche_box.text_frame
tf_fiche.word_wrap = True
p_f = tf_fiche.paragraphs[0]
p_f.alignment = PP_ALIGN.CENTER
r_f = p_f.add_run()
r_f.text = "📊  Fiche Produit KELIA"
r_f.font.size = Pt(10)
r_f.font.bold = True
r_f.font.color.rgb = BLANC
p_f2 = tf_fiche.add_paragraph()
p_f2.alignment = PP_ALIGN.CENTER
r_f2 = p_f2.add_run()
r_f2.text = "\n(.xlsx — Template KELIA\nrempli automatiquement)"
r_f2.font.size = Pt(8)
r_f2.font.color.rgb = GRIS_CLAIR

# ══════════════════════════════════════════════════════════════════════════════
# COLONNE 4 — CONTRÔLES  (x: 9.9 → 13.1)
# ══════════════════════════════════════════════════════════════════════════════
COL4_X = 9.85
COL4_W = 3.28

# Flèche 3 → 4
add_arrow(slide, COL3_X + COL3_W + 0.01, 3.0, COL4_X - 0.01, 3.0, color=ROUGE, width=Pt(3))

# En-tête colonne 4
header4 = add_rounded_rect(slide, COL4_X, 0.95, COL4_W, 0.45, ROUGE)
set_text(header4, "CONTRÔLES & VÉRIFICATIONS", Pt(10), bold=True, color=BLANC)

add_textbox(slide, "Validation automatisée par l'IA",
    COL4_X, 1.43, COL4_W, 0.3, Pt(8), color=GRIS_CLAIR, align=PP_ALIGN.CENTER)

# 3 contrôles
controles = [
    (ROUGE, "Recette",
     "Fiche KELIA générée\n↕\nValeurs attendues\n→ Rapport d'écarts"),
    (RGBColor(0xA0, 0x2E, 0x80), "Conformité contractuelle",
     "Paramétrage KELIA\n↕\nDocument contractuel\n(CG / Notice)\n→ Détection anomalies"),
    (RGBColor(0x8B, 0x45, 0x13), "Non-régression",
     "Paramétrage KELIA V1\n↕\nParamétrage KELIA V2\n→ Rapport des différences"),
]

y_ctrl = 1.78
for couleur, titre_ctrl, desc_ctrl in controles:
    ctrl_box = add_rounded_rect(slide, COL4_X + 0.08, y_ctrl, COL4_W - 0.16, 1.45,
                                 couleur, line_color=BLANC, line_width=Pt(0.75))
    tf_c = ctrl_box.text_frame
    tf_c.word_wrap = True
    pc0 = tf_c.paragraphs[0]
    pc0.alignment = PP_ALIGN.LEFT
    rc0 = pc0.add_run()
    rc0.text = titre_ctrl
    rc0.font.size = Pt(9)
    rc0.font.bold = True
    rc0.font.color.rgb = BLANC
    pc1 = tf_c.add_paragraph()
    pc1.alignment = PP_ALIGN.LEFT
    rc1 = pc1.add_run()
    rc1.text = desc_ctrl
    rc1.font.size = Pt(8)
    rc1.font.color.rgb = GRIS_CLAIR
    y_ctrl += 1.57

# ══════════════════════════════════════════════════════════════════════════════
# BAS DE PAGE
# ══════════════════════════════════════════════════════════════════════════════
footer = add_rect(slide, 0, 7.15, 13.33, 0.35, BLEU_TITRE)
add_textbox(slide, "© Accenture — Confidentiel — Application IA Migration KELIA",
    0.2, 7.18, 12.9, 0.28, Pt(8), color=GRIS_CLAIR, align=PP_ALIGN.CENTER)

# ── Sauvegarde ───────────────────────────────────────────────────────────────
output_path = r"C:\Users\walid.ben.lamine\OneDrive - Accenture\01ApplicactionCartoProduit\Schema_Application_KELIA.pptx"
prs.save(output_path)
print(f"OK - Fichier genere : {output_path}")
